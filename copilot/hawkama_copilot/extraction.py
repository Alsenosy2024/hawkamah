"""Deterministic file-text extraction for the ingestion pipeline.

Ported from the JS `fileExtraction.ts`: extract *real* text from each upload, never
fabricate. Native parsers for the structured formats (DOCX/PPTX/XLSX/PDF/text);
Gemini multimodal as the fallback for scanned PDFs and images. The result carries
the method used and a human reason when extraction yields nothing, so the agent
can tell the operator "this file was an image with no readable text" instead of
silently dropping it.
"""

from __future__ import annotations

import io
import mimetypes
from dataclasses import dataclass
from pathlib import Path

from google.genai import types

from . import genai_client
from .config import MODELS

# Bytes over which we route binary docs to the Files API rather than inlining.
_INLINE_LIMIT = 14 * 1024 * 1024

_TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".xml", ".log", ".rtf"}
_HTML_EXT = {".html", ".htm"}
_DOCX_EXT = {".docx"}
_PPTX_EXT = {".pptx"}
_XLSX_EXT = {".xlsx", ".xlsm"}
_PDF_EXT = {".pdf"}
_IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}


@dataclass
class ExtractResult:
    text: str
    method: str          # text|html|docx|pptx|xlsx|pdf|gemini|none
    error: str | None = None

    @property
    def ok(self) -> bool:
        return bool(self.text.strip())


def extract_file(path: str | Path) -> ExtractResult:
    """Extract text from a file on disk by extension/sniffing."""
    p = Path(path)
    data = p.read_bytes()
    return extract_bytes(data, filename=p.name)


def is_image(filename: str) -> bool:
    return Path(filename).suffix.lower() in _IMG_EXT


def caption_image(data: bytes, filename: str) -> str:
    """Short Arabic caption + any visible text for an image, used as the readable
    evidence/citation text alongside the image's multimodal embedding."""
    mime = mimetypes.guess_type(filename)[0] or "image/png"
    prompt = (
        "صف محتوى هذه الصورة بإيجاز (سطر إلى ثلاثة أسطر) بالعربية، واذكر أي نص ظاهر فيها "
        "حرفيًا إن وجد. لا تضف مقدمات."
    )
    try:
        client = genai_client.get_client()
        resp = client.models.generate_content(
            model=MODELS.text,
            contents=[prompt, types.Part.from_bytes(data=data, mime_type=mime)],
        )
        return (resp.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def extract_bytes(data: bytes, *, filename: str) -> ExtractResult:
    """Extract text from in-memory bytes (used by the upload API)."""
    ext = Path(filename).suffix.lower()
    try:
        if ext in _TEXT_EXT:
            return ExtractResult(_decode(data), "text")
        if ext in _HTML_EXT:
            return ExtractResult(_html_to_text(data), "html")
        if ext in _DOCX_EXT:
            return ExtractResult(_docx_to_text(data), "docx")
        if ext in _PPTX_EXT:
            return ExtractResult(_pptx_to_text(data), "pptx")
        if ext in _XLSX_EXT:
            return ExtractResult(_xlsx_to_text(data), "xlsx")
        if ext in _PDF_EXT:
            return _pdf_to_text(data, filename)
        if ext in _IMG_EXT:
            return _gemini_extract(data, filename, "image")
        # Unknown: try utf-8, else last-resort Gemini.
        txt = _decode(data)
        if txt.strip():
            return ExtractResult(txt, "text")
        return _gemini_extract(data, filename, "binary")
    except Exception as e:  # noqa: BLE001 — never let one bad file crash ingest
        return ExtractResult("", "none", error=f"{type(e).__name__}: {e}")


# --------------------------------------------------------------------------- #
# Native parsers                                                               #
# --------------------------------------------------------------------------- #
def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "cp1256", "latin-1"):  # cp1256 = Arabic Windows
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _html_to_text(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(data), "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())


def _docx_to_text(data: bytes) -> str:
    import docx  # python-docx

    doc = docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            # Preserve heading semantics so chunking can use them.
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading"):
                level = "".join(c for c in style if c.isdigit()) or "2"
                parts.append("#" * min(int(level), 6) + " " + para.text.strip())
            else:
                parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx_to_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)


def _xlsx_to_text(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(",".join(cells))
    wb.close()
    return "\n".join(parts)


def _pdf_to_text(data: bytes, filename: str) -> ExtractResult:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for n, page in enumerate(reader.pages, 1):
            txt = (page.extract_text() or "").strip()
            if txt:
                pages.append(f"<!-- page {n} -->\n{txt}")
        joined = "\n\n".join(pages)
        # Scanned PDFs yield little/no embedded text → multimodal OCR fallback.
        if len(joined.strip()) >= 40:
            return ExtractResult(joined, "pdf")
    except Exception:  # noqa: BLE001 — fall through to Gemini
        pass
    return _gemini_extract(data, filename, "pdf")


# --------------------------------------------------------------------------- #
# Gemini multimodal fallback (scanned PDFs, images)                           #
# --------------------------------------------------------------------------- #
def _gemini_extract(data: bytes, filename: str, kind: str) -> ExtractResult:
    if len(data) > _INLINE_LIMIT * 4:
        return ExtractResult("", "none", error=f"{kind} too large to OCR ({len(data)} bytes)")
    mime = mimetypes.guess_type(filename)[0] or ("application/pdf" if kind == "pdf" else "image/png")
    prompt = (
        "استخرج كل النص الظاهر في هذا الملف حرفيًا، محافظًا على الترتيب والعناوين "
        "والجداول قدر الإمكان. لا تضف أي شرح أو تعليق. إن لم يوجد نص مقروء فأعد سطرًا "
        "فارغًا فقط.\n"
        "Extract all visible text verbatim, preserving order, headings and tables. "
        "Do not add commentary. If there is no readable text, return an empty line."
    )
    try:
        client = genai_client.get_client()
        if len(data) > _INLINE_LIMIT:
            up = client.files.upload(file=io.BytesIO(data), config={"mime_type": mime})
            contents = [prompt, up]
        else:
            contents = [prompt, types.Part.from_bytes(data=data, mime_type=mime)]
        resp = client.models.generate_content(model=MODELS.text, contents=contents)
        text = (resp.text or "").strip()
        if text:
            return ExtractResult(text, "gemini")
        return ExtractResult("", "none", error=f"no readable text in {kind}")
    except Exception as e:  # noqa: BLE001
        return ExtractResult("", "none", error=f"gemini OCR failed: {type(e).__name__}: {e}")
