"""PPTX export (python-pptx) with RTL paragraphs.

Turns a generated document into a deck: a cover slide, a divider slide per H2,
and content slides that page long bodies (bullets chunked) and render tables.
python-pptx needs the `<a:pPr algn="r" rtl="1">` bit set on each paragraph for
correct Arabic direction, which we do via the lxml element.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .markdown_ast import Block, parse_markdown

NAVY = RGBColor(0x1A, 0x35, 0x57)
GOLD = RGBColor(0xC8, 0x91, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x29, 0x37)
MIST = RGBColor(0xEE, 0xF2, 0xFF)
FONT = "Almarai"
MAX_BULLETS = 8


def export_pptx(markdown: str, title: str, *, company: str = "") -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _cover(prs, blank, title, company)

    blocks = parse_markdown(markdown)
    i = 0
    n = len(blocks)
    while i < n:
        b = blocks[i]
        if b.type == "heading" and b.level <= 2:
            _divider(prs, blank, b.text)
            i += 1
            # gather following content until next H2
            buf: list[Block] = []
            while i < n and not (blocks[i].type == "heading" and blocks[i].level <= 2):
                buf.append(blocks[i])
                i += 1
            _content_slides(prs, blank, b.text, buf)
        else:
            i += 1
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _rtl(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("algn", "r")
    pPr.set("rtl", "1")


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _cover(prs, layout, title: str, company: str) -> None:
    s = prs.slides.add_slide(layout)
    _bg(s, NAVY)
    tf = _textbox(s, 0.8, 2.6, 11.7, 2.4)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    _rtl(p)
    if company:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = company
        r2.font.size = Pt(22)
        r2.font.color.rgb = GOLD
        r2.font.name = FONT
        _rtl(p2)


def _divider(prs, layout, text: str) -> None:
    s = prs.slides.add_slide(layout)
    _bg(s, MIST)
    tf = _textbox(s, 0.8, 3.0, 11.7, 1.6)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = FONT
    _rtl(p)


def _content_slides(prs, layout, heading: str, blocks: list[Block]) -> None:
    bullets: list[str] = []

    def flush_bullets() -> None:
        nonlocal bullets
        for i in range(0, len(bullets), MAX_BULLETS):
            _bullet_slide(prs, layout, heading, bullets[i : i + MAX_BULLETS])
        bullets = []

    for b in blocks:
        if b.type in ("bullet", "ordered", "paragraph", "quote"):
            bullets.append(b.text)
        elif b.type == "heading":
            flush_bullets()
            # sub-heading becomes its own bullet header
            bullets.append("▸ " + b.text)
        elif b.type == "table":
            flush_bullets()
            _table_slide(prs, layout, heading, b)
    flush_bullets()


def _bullet_slide(prs, layout, heading: str, items: list[str]) -> None:
    s = prs.slides.add_slide(layout)
    _bg(s, WHITE)
    head = _textbox(s, 0.6, 0.4, 12.1, 0.9)
    hp = head.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.size = Pt(24)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    hr.font.name = FONT
    _rtl(hp)

    body = _textbox(s, 0.6, 1.4, 12.1, 5.6)
    first = True
    for it in items:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        r = p.add_run()
        r.text = "• " + it[:300]
        r.font.size = Pt(16)
        r.font.color.rgb = INK
        r.font.name = FONT
        _rtl(p)


def _table_slide(prs, layout, heading: str, b: Block) -> None:
    s = prs.slides.add_slide(layout)
    _bg(s, WHITE)
    head = _textbox(s, 0.6, 0.3, 12.1, 0.8)
    hp = head.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.size = Pt(22)
    hr.font.bold = True
    hr.font.color.rgb = NAVY
    hr.font.name = FONT
    _rtl(hp)

    rows = min(len(b.rows) + 1, 12)
    cols = max(len(b.headers), 1)
    shape = s.shapes.add_table(rows, cols, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8))
    table = shape.table
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = b.headers[c] if c < len(b.headers) else ""
        _cell_style(cell, header=True)
    for ri, row in enumerate(b.rows[: rows - 1], start=1):
        for c in range(cols):
            cell = table.cell(ri, c)
            cell.text = row[c] if c < len(row) else ""
            _cell_style(cell, header=False)


def _cell_style(cell, *, header: bool) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY if header else WHITE
    for p in cell.text_frame.paragraphs:
        _rtl(p)
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.name = FONT
            r.font.color.rgb = WHITE if header else INK
            r.font.bold = header
